import os
import json
import requests
from urllib.parse import quote
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from io import BytesIO

# Load .env from the directory containing tools.py
env_path = os.path.join(os.path.dirname(__file__), '.env')
load_dotenv(dotenv_path=env_path)

def get_llm_client_and_model():
    """Returns an initialized OpenAI client and model name based on available env variables."""
    if os.getenv("GEMINI_API_KEY"):
        return OpenAI(
            api_key=os.getenv("GEMINI_API_KEY"),
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
        ), "gemini-2.5-flash"
    elif os.getenv("XAI_API_KEY"):
        return OpenAI(
            api_key=os.getenv("XAI_API_KEY"),
            base_url="https://api.x.ai/v1"
        ), "grok-beta"
    else:
        return OpenAI(api_key=os.getenv("OPENAI_API_KEY")), "gpt-4o-mini"

def generate_slide_content(topic: str, subtopic: str) -> dict:
    """
    Generates slide content including title and bullet points.

    Args:
        topic (str): Main topic of presentation.
        subtopic (str): Specific section of the topic.

    Returns:
        dict: {
            "title": str,
            "bullets": List[str]
        }
    """
    client, model = get_llm_client_and_model()
    
    prompt = f"""
    You are an expert presentation creator. Provide content for a single slide.
    Main Topic: {topic}
    Slide Subtopic: {subtopic}
    
    Output strictly a JSON object with this exact structure, containing the Slide title and exactly 3-5 concise bullet points:
    {{
        "title": "Clear Slide Title based on Subtopic",
        "bullets": ["Point 1", "Point 2", "Point 3"]
    }}
    """
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        content = response.choices[0].message.content
        return json.loads(content)
    except Exception as e:
        print(f"Error generating slide content for {subtopic}: {e}")
        return {"title": subtopic, "bullets": ["Error generating content. Please verify your OpenAI API Key."]}

def fetch_image(query: str) -> str:
    """
    Fetches a relevant image URL based on query.

    Args:
        query (str): Description of image.

    Returns:
        str: Image URL.
    """
    encoded_query = quote(query)
    # Using pollinations.ai for free, prompt-based image generation
    url = f"https://image.pollinations.ai/prompt/{encoded_query}?width=800&height=600&nologo=true"
    return url

def create_slide(title: str, bullets: list, image_url: str = None) -> dict:
    """
    Creates a structured slide object.

    Args:
        title (str): Slide title.
        bullets (list): Bullet points.
        image_url (str, optional): Image link.

    Returns:
        dict: Structured slide.
    """
    return {
        "title": title,
        "bullets": bullets,
        "image": image_url
    }

def build_ppt(slides: list, file_name: str) -> str:
    """
    Converts slides into a PPT file.

    Args:
        slides (list): List of slide dictionaries.
        file_name (str): Output file name.

    Returns:
        str: File path of generated PPT.
    """
    prs = Presentation()
    
    for slide_data in slides:
        title_str = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        image_url = slide_data.get("image")
        
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = title_str
        tf = body.text_frame
        if bullets:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        
        if image_url:
            # Adjust body text width to make room for image on the right
            body.width = Inches(4.5)
            try:
                response = requests.get(image_url, timeout=10)
                if response.status_code == 200:
                    image_stream = BytesIO(response.content)
                    slide.shapes.add_picture(
                        image_stream, 
                        Inches(5.0), # Left
                        Inches(2.0), # Top
                        width=Inches(4.5)
                    )
            except Exception as e:
                print(f"Failed to fetch image for slide '{title_str}': {e}")
                
    prs.save(file_name)
    return os.path.abspath(file_name)
